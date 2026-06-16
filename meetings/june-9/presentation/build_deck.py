"""Build the mentor progress deck (PowerPoint) from the simulation figures.

Usage: python3 build_deck.py   (run from the presentation folder)
Produces progress-deck.pptx in the same folder; reads figures from the sibling
deliverables folder.

Dark theme with a red accent. One title slide plus five content slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
# Figures live in the sibling deliverables folder; the pptx is written next to
# this script (the presentation folder).
FIG = os.path.normpath(os.path.join(HERE, "..", "deliverables"))

DARK = RGBColor(0x1E, 0x1E, 0x26)
LIGHT = RGBColor(0xEC, 0xEC, 0xEC)
MUTED = RGBColor(0xB2, 0xB2, 0xBE)
RED = RGBColor(0xE0, 0x3A, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK


def box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def line(tf, text, size, color, bold=False, first=False, space=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"
    return p


def title_of(slide, text):
    tf = box(slide, 0.7, 0.45, 12.0, 1.1)
    line(tf, text, 30, RED, bold=True, first=True)


def bullets_of(slide, items, left=0.75, top=1.7, width=12.0, size=19):
    tf = box(slide, left, top, width, 5.3)
    for i, b in enumerate(items):
        line(tf, "•  " + b, size, LIGHT, first=(i == 0), space=12)


def caption(slide, text, left, top, width):
    tf = box(slide, left, top, width, 0.4)
    line(tf, text, 12, MUTED, first=True)


def img(slide, name, left, top, width):
    slide.shapes.add_picture(os.path.join(FIG, name), Inches(left), Inches(top), width=Inches(width))


# Slide 1: title
s = prs.slides.add_slide(BLANK)
add_bg(s)
tf = box(s, 0.9, 2.0, 11.5, 2.6)
line(tf, "Learning Emerging Behaviors of Dynamic Multi-Agent Systems "
         "using Graph Neural Networks", 30, LIGHT, bold=True, first=True, space=14)
line(tf, "Week 1 progress: a boids flocking simulation", 20, RED)
tf = box(s, 0.9, 5.0, 11.5, 1.6)
line(tf, "Mustafa Nazeer", 22, LIGHT, bold=True, first=True, space=6)
line(tf, "Mentor: Dr. Diego Patino", 16, MUTED, space=4)
line(tf, "June 16, 2026", 16, MUTED)
tf = box(s, 0.9, 6.95, 11.5, 0.4)
line(tf, "Figures and code produced with AI assistance, disclosed per course policy.",
     11, MUTED, first=True)

# Slide 2: problem and background
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Emergent behavior from local rules")
bullets_of(s, [
    "Many agent systems (flocks, herds, crowds, traffic) show group behavior "
    "that no single agent is told to produce.",
    "Reynolds' boids (1987): three local rules produce flocking, separation, "
    "alignment, and cohesion.",
    "Project goal: a Graph Neural Network that learns these interaction rules "
    "from agent trajectories.",
    "This week: build the rule based simulation that generates that behavior "
    "and the data the network would learn from.",
])

# Slide 3: what we built
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "A boids simulation with our own rules")
bullets_of(s, [
    "Pure Python (numpy, scipy, matplotlib); 120 boids in a 60 by 60 "
    "wraparound world.",
    "The classic three rules plus avoidance of circular obstacles.",
    "Custom rule: a predator that triggers collective escape and splitting.",
    "Inverse distance separation keeps the flock spaced out instead of "
    "collapsing to a point.",
    "Reproducible, with a unit tested core of 24 passing tests.",
])

# Slide 4: predator escape demo
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Collective escape from a predator")
bullets_of(s, [
    "The flock forms and flows around obstacles as one coordinated group.",
    "A predator (red star) chases the nearest boids; nearby boids flee.",
    "The flock splits, scatters, then streams back together while fleeing.",
], top=1.55, size=18)
img(s, "snapshot_early.png", 1.7, 3.5, 3.6)
img(s, "snapshot_late.png", 8.0, 3.5, 3.6)
caption(s, "Flocking and obstacle avoidance", 1.7, 7.05, 3.6)
caption(s, "Splitting and escape", 8.0, 7.05, 3.6)
caption(s, "Animated version: flock_escape.gif", 5.4, 3.0, 4.5)

# Slide 5: results and findings (no graphs, general statements)
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Results and findings")
bullets_of(s, [
    "From random starts, the boids organize themselves into one coherent, "
    "aligned flock within roughly 40 frames.",
    "At rest the flock moves almost perfectly in unison as a single group, with "
    "clear spacing between individuals rather than collapsing to a point.",
    "The flock reliably flows around the obstacles without collisions.",
    "The predator drives a clear collective response: the flock splits, "
    "alignment breaks down, and the boids scatter as they flee, then stream "
    "back together once it passes.",
    "These findings are measurable, not just visual, which is what makes the "
    "behavior learnable by the planned Graph Neural Network.",
])

# Slide 6: next steps
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Next steps")
bullets_of(s, [
    "Port the model to NVIDIA Isaac Sim on a GPU machine, and add a 3D "
    "environment.",
    "Validate against real flocking datasets.",
    "Build the Graph Neural Network that learns these interaction rules from "
    "trajectories.",
    "Note: built in Python first because the development machine has no "
    "NVIDIA GPU.",
])

out = os.path.join(HERE, "progress-deck.pptx")
prs.save(out)
print("Wrote", out)
