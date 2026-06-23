"""Build the mentor progress deck (PowerPoint) from the simulation figures.

Usage: python3 build_deck.py   (run from the presentation folder)
Produces progress-deck.pptx in the same folder; reads figures from the sibling
deliverables folder.

Dark theme with a red accent. One title slide plus five content slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
# Figures live in the sibling deliverables folder; the pptx is written next to
# this script (the presentation folder).
FIG = os.path.normpath(os.path.join(HERE, "..", "deliverables"))

DARK = RGBColor(0x1E, 0x1E, 0x26)
LIGHT = RGBColor(0xEC, 0xEC, 0xEC)
MUTED = RGBColor(0xB2, 0xB2, 0xBE)
RED = RGBColor(0xE0, 0x3A, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK


def box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def line(tf, text, size, color, bold=False, first=False, space=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"
    return p


def title_of(slide, text):
    tf = box(slide, 0.7, 0.45, 12.0, 1.1)
    line(tf, text, 30, RED, bold=True, first=True)


def bullets_of(slide, items, left=0.75, top=1.7, width=12.0, size=19):
    tf = box(slide, left, top, width, 5.3)
    for i, b in enumerate(items):
        line(tf, "•  " + b, size, LIGHT, first=(i == 0), space=12)


def caption(slide, text, left, top, width):
    tf = box(slide, left, top, width, 0.4)
    line(tf, text, 12, MUTED, first=True)


def img(slide, name, left, top, width):
    slide.shapes.add_picture(os.path.join(FIG, name), Inches(left), Inches(top), width=Inches(width))


# Slide 1: title
s = prs.slides.add_slide(BLANK)
add_bg(s)
tf = box(s, 0.9, 2.0, 11.5, 2.6)
line(tf, "Learning Emerging Behaviors of Dynamic Multi-Agent Systems "
         "using Graph Neural Networks", 30, LIGHT, bold=True, first=True, space=14)
line(tf, "Progress: a trajectory data pipeline for the GNN", 20, RED)
tf = box(s, 0.9, 5.0, 11.5, 1.6)
line(tf, "Mustafa Nazeer and Joel Pampam", 22, LIGHT, bold=True, first=True, space=6)
line(tf, "Mentor: Dr. Diego Patino", 16, MUTED, space=4)
line(tf, "June 23, 2026", 16, MUTED)
tf = box(s, 0.9, 6.95, 11.5, 0.4)
line(tf, "Figures and code produced with AI assistance, disclosed per course policy.",
     11, MUTED, first=True)

# Slide 2: this week's goal
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "From a visual flock to usable data")
bullets_of(s, [
    "Last meeting set the next step: stop treating the simulation as only a "
    "picture, and start grabbing data from it.",
    "At every step we capture each agent's position, velocity, and who its "
    "neighbors are, then write it all to a CSV.",
    "That CSV is the trajectory data the Graph Neural Network will later learn "
    "from.",
    "We aligned our output with Angel Solis' reference pipeline so the two are "
    "interchangeable.",
])

# Slide 3: the data pipeline
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "What the recorder captures")
bullets_of(s, [
    "One row per agent per step: step, id, position, velocity.",
    "A pre_planned flag marking leader agents that follow set paths.",
    "Three neighbor lists per agent, one each for separation, cohesion, and "
    "alignment ranges.",
    "Those neighbor lists are the time varying interaction graph: who is "
    "influencing whom, at every instant.",
    "The schema matches Angel's exactly, so our pure Python data and his feed "
    "the same model.",
])

# Slide 4: pre-planned debug agents
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Known paths to trust the data")
bullets_of(s, [
    "Pre-planned agents follow exact shapes (figure 8 and straight lines) we "
    "know in advance.",
    "Run alone, they let us open the CSV and confirm the logged path matches "
    "the real one.",
    "Dropped into the flock as leaders (gold), the normal boids react to them.",
], top=1.55, size=18)
img(s, "snapshot_preplanned.png", 4.6, 3.2, 4.1)
caption(s, "Pre-planned agents tracing known figure 8 and line paths", 4.6, 7.05, 4.5)

# Slide 5: results and findings (no graphs, qualitative)
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Results and findings")
bullets_of(s, [
    "The simulation now produces a clean CSV of full trajectories for every "
    "agent across the whole run.",
    "Reflective walls keep each path continuous, with no sudden jumps, which "
    "matters for training a model on the motion.",
    "The neighbor lists capture how the interaction graph changes over time as "
    "the flock moves and regroups.",
    "The known path agents confirm the recorded data faithfully matches the "
    "real motion, so we can trust it.",
    "We now have the raw material a Graph Neural Network would consume.",
])

# Slide 6: next steps
s = prs.slides.add_slide(BLANK)
add_bg(s)
title_of(s, "Next steps")
bullets_of(s, [
    "Port the model to NVIDIA Isaac Sim on a GPU machine, using Angel's repo "
    "as the reference.",
    "Feed this trajectory CSV into a first Graph Neural Network experiment.",
    "Confirm the intended maximum speed with Angel (the notes say 5, his config "
    "says 10).",
    "Read two to three recent papers on learning multi-agent interactions.",
])

out = os.path.join(HERE, "progress-deck.pptx")
prs.save(out)
print("Wrote", out)
